from pptx import Presentation


# создаем новую презентацию
prs = Presentation()
# получаем схему расположения элементов для заголовочного слайда
title_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Python random"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = "Модуль random предоставляет функции для генерации случайных чисел, " \
                "букв, случайного выбора элементов последовательности."

slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Shuffle"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = "Метод random.shuffle() используется для перемешивания" \
                " последовательности.\nМетод shuffle() смешивает элементы списка на месте."

slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Randint"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = 'Random.randint(a, b) - случайное целое число от a до b\n' \
                'работает только с целыми числами;\nПринимает только два параметра;'
slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Choice"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = 'random.choice(sequence) - передает случайный элемент последовательности.\n'\
                'Введен в версии Python 3.6.\nПозволяет ' \
                'повторить несколько раз один элемент.'

slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sample"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = 'Метод random.sample() выбирает' \
                ' несколько элементов из последовательности.\n' \
                'Возвращает список уникальных элементов длинной k, выбранных из последова' \
                'тельности'

slide = prs.slides.add_slide(title_slide_layout)
# создаем у слайда заголовок и текст
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Uniform"
tf = subtitle.text_frame
p = tf.add_paragraph()
p.font.name = 'Courier New'
p.text = 'random.uniform() используется для генерации числа с плавающей запятой в' \
                ' заданном промежутка;\nЗначение конечной точки может включаться' \
                ' в диапазон, но это не обязательно.'

prs.save('tet.pptx')
